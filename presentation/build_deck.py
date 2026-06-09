# -*- coding: utf-8 -*-
"""Build the Creative Bowl Lab final-presentation draft (15 slides, animated)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import deck_kit as k
import animations as anim

ASSETS = r"C:\projects\CBL\presentation\assets"
OUT = r"C:\projects\CBL\presentation\Creative_Bowl_Lab_Final_DRAFT.pptx"

prs = Presentation()
prs.slide_width = Inches(k.SW)
prs.slide_height = Inches(k.SH)
BLANK = prs.slide_layouts[6]
M = k.MARGIN
CW = k.CW


def img(name):
    return os.path.join(ASSETS, name)


def E(shape, type="fade", trigger="after", dur=600, delay=120):
    return {"shape": shape, "type": type, "trigger": trigger, "dur": dur, "delay": delay}


def head(slide, kick, title, accent=k.GOLD, tsize=34):
    kk, tt, rl = k.header(slide, kick, title, accent=accent, title_size=tsize)
    return [E(kk, "fade", "after", 460, 120),
            E(tt, "float", "after", 640, 40),
            E(rl, "wipe", "after", 420, 0)]


# ============================================================ 01 · TITLE
def s01():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    g1 = k.glow(s, 3.1, 3.0, 4.6, k.VIOLET, inner_alpha=34)
    g2 = k.glow(s, 10.6, 4.4, 4.4, k.CYAN, inner_alpha=26)
    kick = k.kicker(s, M, 1.18, "TU/e  ·  MULTI-DISCIPLINARY CBL  ·  GROUP 5", color=k.GOLD)
    t1 = k.text(s, M, 1.62, CW, 1.7, "Creative Bowl Lab", font=k.DISPLAY,
                size=72, color=k.INK, line_spacing=0.98)
    rl = k.rule(s, M + 0.04, 3.16, 2.0, color=k.GOLD, weight=2.4)
    sub = k.text(s, M, 3.40, 10.6, 1.2,
                 "Making meditation visible — where a singing bowl, a heartbeat,\nand a living body become one image.",
                 font=k.DISPLAY, italic=True, size=23, color=k.BODY, line_spacing=1.12)
    bar = k.chakra_bar(s, M, 5.62, 6.2, h=0.16)
    when = k.text(s, M, 6.02, 8, 0.4, "FINAL PRESENTATION  ·  19 JUNE 2026",
                  font=k.SANS, size=11.5, color=k.GOLD, bold=True, tracking=2.6)
    names = k.text(s, M, 6.46, CW, 0.5,
                   "Arav Madan   ·   Joris Geleijns   ·   Henk   ·   Alice Ribeiro Mota Vieira   ·   Mahiraa   ·   Alejandra Otero Bozzano",
                   font=k.SANS, size=11, color=k.MUTE, tracking=0.4)
    eff = [E(g1, "fade", "after", 900, 0), E(g2, "fade", "with", 900, 0),
           E(kick, "fade", "after", 460, 120), E(t1, "float", "after", 760, 60),
           E(rl, "wipe", "after", 460, 0), E(sub, "fade", "after", 620, 60)]
    eff += [E(seg, "wipe", "with" if i else "after", 360, 0) for i, seg in enumerate(bar)]
    eff += [E(when, "fade", "after", 460, 80), E(names, "fade", "after", 520, 40)]
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 800)


# ============================================================ 02 · TEAM
def s02():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 11.5, 1.2, 3.4, k.VIOLET, 18)
    eff = head(s, "SIX DISCIPLINES, ONE INSTALLATION", "A living artwork, built across the table")
    team = [
        ("COMPUTER SCIENCE", "Arav Madan", "Signal pipeline, TouchDesigner, Arduino link", k.CYAN),
        ("PSYCHOLOGY & TECH", "Joris Geleijns", "Meaning, perception, the health research", k.VIOLET),
        ("ARCHITECTURE", "Henk", "Space, the participant's position, atmosphere", k.GOLD),
        ("CHEMICAL ENG.", "Alice R. M. Vieira", "Bowl frequency experiments, materials", k.GREEN),
        ("INDUSTRIAL DESIGN", "Mahiraa", "Requirements, Arduino & physical build", k.ORANGE),
        ("MECHANICAL ENG.", "Alejandra O. Bozzano", "Cymatics model, report architecture", k.ROSE),
    ]
    cw, ch, gx, gy = 3.62, 1.78, 0.30, 0.32
    x0, y0 = M, 2.18
    for i, (disc, name, role, col) in enumerate(team):
        cx = x0 + (i % 3) * (cw + gx)
        cy = y0 + (i // 3) * (ch + gy)
        c = k.card(s, cx, cy, cw, ch, fill=k.PANEL, stroke=k.HAIR, radius=0.07)
        tab = k.box(s, cx, cy + 0.18, 0.06, ch - 0.36, MSO_SHAPE.ROUNDED_RECTANGLE)
        k.solid(tab, col)
        k.text(s, cx + 0.30, cy + 0.24, cw - 0.5, 0.3, disc, font=k.SANS, size=10.5,
               color=col, bold=True, tracking=1.8)
        k.text(s, cx + 0.30, cy + 0.58, cw - 0.5, 0.4, name, font=k.DISPLAY, size=21, color=k.INK)
        k.text(s, cx + 0.30, cy + 1.04, cw - 0.5, 0.6, role, font=k.SANS, size=11.5,
               color=k.MUTE, line_spacing=1.1)
        eff.append(E(c, "float", "after" if i == 0 else ("with" if i % 3 else "after"), 520, 60))
    k.footer(s, 2, section="The team")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 03 · GOAL
def s03():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 2.4, 5.6, 4.2, k.VIOLET, 22)
    eff = head(s, "THE GOAL", "Can you watch yourself calm down?")
    body = k.text(s, M, 2.32, 6.5, 2.6,
                  "Students live under near-constant stress. Meditation and sound can genuinely settle the body — but that shift is invisible. You feel it; you never see it.\n\nOur goal: an interactive installation that turns a moment of calm into something you can watch happen, in real time, around your own body.",
                  font=k.SANS, size=16.5, color=k.BODY, line_spacing=1.28)
    # "why this goal is ambitious" card — speaks to Crit 1 (challenging/original/innovative)
    qc = k.card(s, 8.0, 2.32, 4.42, 3.7, fill=k.PANEL2, stroke=k.HAIR, radius=0.06)
    k.text(s, 8.34, 2.62, 3.8, 0.34, "WHY IT'S A CHALLENGING GOAL", font=k.SANS,
           size=11, color=k.GOLD, bold=True, tracking=2.0)
    pts = [("Original.", "There is no off-the-shelf version of this — we designed the whole sound-to-image chain ourselves.", k.CYAN),
           ("Complex.", "Fusing live bowl sound, a heartbeat and body-tracking into one GPU image, in real time.", k.VIOLET),
           ("Innovative.", "It makes an invisible state — your body calming — into something you can watch.", k.GREEN)]
    py = 3.06
    pshapes = []
    for term, bodytxt, col in pts:
        r = k.rich(s, 8.34, py, 3.8, 0.9, [(term + "  ", col, True), (bodytxt, k.BODY, False)],
                   size=12.5, line_spacing=1.16)
        pshapes.append(r)
        py += 0.92
    eff += [E(body, "fade", "after", 640, 80), E(qc, "float", "after", 560, 120)]
    eff += [E(r, "fade", "after" if i == 0 else "after", 380, 40) for i, r in enumerate(pshapes)]
    k.footer(s, section="Goal & motivation")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 04 · SCIENCE
def s04():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    eff = head(s, "WHY IT WORKS — THE EVIDENCE", "Sound slows the body down")
    cards = [
        ("↑ HRV", "Singing bowls", "Bowl meditation lowered the stress index and raised heart-rate variability versus silence.", "Trivedi & Saboo, 2019", k.GREEN),
        ("↓ HR / BP", "Music & the heart", "Music increases parasympathetic (“rest-and-digest”) activity and lowers heart rate and blood pressure.", "Mojtabavi et al., 2020", k.CYAN),
        ("↓ Anxiety", "Bowl therapy", "A 2025 systematic review links singing-bowl therapy to reduced anxiety and improved sleep.", "Cai et al., 2025", k.VIOLET),
    ]
    cw, gx, x0, y0, ch = 3.78, 0.30, M, 2.32, 3.0
    for i, (stat, title, body, cite, col) in enumerate(cards):
        cx = x0 + i * (cw + gx)
        c = k.card(s, cx, y0, cw, ch, fill=k.PANEL, stroke=k.HAIR, radius=0.06)
        k.text(s, cx + 0.34, y0 + 0.30, cw - 0.6, 0.6, stat, font=k.DISPLAY, size=30, color=col)
        k.text(s, cx + 0.34, y0 + 1.00, cw - 0.6, 0.4, title, font=k.SANS, size=14.5,
               color=k.INK, bold=True, tracking=0.4)
        k.text(s, cx + 0.34, y0 + 1.46, cw - 0.66, 1.2, body, font=k.SANS, size=12.5,
               color=k.BODY, line_spacing=1.2)
        k.text(s, cx + 0.34, y0 + 2.58, cw - 0.6, 0.3, cite, font=k.SANS, size=10.5,
               color=k.MUTE, italic=True)
        eff.append(E(c, "float", "after" if i == 0 else "with", 560, 80))
    strip = k.card(s, M, 5.66, CW, 0.66, fill="120F18", stroke=k.AMBER, radius=0.18, stroke_w=1.2)
    k.rich(s, M + 0.34, 5.66, CW - 0.7, 0.66,
           [("An honest framing:  ", k.AMBER, True),
            ("we present this as relaxation and mood support — a calming experience — not a proven medical treatment.", k.BODY, False)],
           size=13.5, anchor=MSO_ANCHOR.MIDDLE)
    eff.append(E(strip, "fade", "after", 520, 80))
    k.footer(s, 4, section="Scientific grounding")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 05 · TRIPLE DIAMOND
def s05():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 6.7, 6.4, 4.6, k.VIOLET, 16)
    eff = head(s, "HOW WE GOT HERE", "A triple-diamond design process")
    diamonds = [
        ("01", "EXPLORE", "Photography, music,\npoetry, engineering art", k.CYAN),
        ("02", "CONVERGE", "Bowl + sound,  lens filter,\npoetry with sensors", k.VIOLET),
        ("03", "REFOCUS", "One clear story:  bowl +\nheartbeat + body", k.GOLD),
    ]
    dsize = 2.2
    gap = (CW - 3 * dsize) / 2
    y = 2.7
    for i, (num, label, sub, col) in enumerate(diamonds):
        cx = M + i * (dsize + gap)
        d = k.box(s, cx, y, dsize, dsize, MSO_SHAPE.DIAMOND)
        k.solid(d, k.PANEL2)
        k.line(d, col, 1.6)
        k.text(s, cx, y + 0.62, dsize, 0.5, num, font=k.DISPLAY, size=30, color=col,
               align=PP_ALIGN.CENTER)
        k.text(s, cx, y + 1.16, dsize, 0.4, label, font=k.SANS, size=12.5, color=k.INK,
               bold=True, tracking=2.2, align=PP_ALIGN.CENTER)
        cap = k.text(s, cx - 0.2, y + dsize + 0.12, dsize + 0.4, 0.9, sub, font=k.SANS,
                     size=12, color=k.MUTE, align=PP_ALIGN.CENTER, line_spacing=1.18)
        if i < 2:
            ar = k.text(s, cx + dsize + gap / 2 - 0.3, y + 0.7, 0.6, 0.8, "›",
                        font=k.SANS, size=40, color=k.HAIR, align=PP_ALIGN.CENTER)
            eff.append(E(ar, "fade", "with", 300, 0))
        eff.append(E(d, "zoom", "after" if i == 0 else "after", 560, 60))
        eff.append(E(cap, "fade", "with", 460, 0))
    note = k.text(s, M, 6.46, CW, 0.4,
                  "Each diamond widened with ideas, then narrowed with research and a €100 reality check.",
                  font=k.DISPLAY, italic=True, size=15, color=k.BODY, align=PP_ALIGN.CENTER)
    eff.append(E(note, "fade", "after", 460, 60))
    k.footer(s, 5, section="Design process")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 06 · DISCARDED
def s06():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    eff = head(s, "DECISIONS & MID-TERM FEEDBACK", "Why the scope is what it is")
    fb = k.card(s, M, 2.02, CW, 0.6, fill="120F18", stroke=k.ROSE, radius=0.16, stroke_w=1.2)
    k.rich(s, M + 0.34, 2.02, CW - 0.7, 0.6,
           [("Mid-term jury:  ", k.ROSE, True),
            ("“not visualised · too complex · unclear.”   Our response — cut to one clear, visible story. Every decision below follows from that.", k.BODY, False)],
           size=12.5, anchor=MSO_ANCHOR.MIDDLE)
    eff.append(E(fb, "fade", "after", 460, 60))
    rows = [
        ("Photography", "Dropped", "Too static — it never read as a living, reacting thing.", k.MUTE, False),
        ("Lens / camera filter", "Dropped", "Expensive and physically large for an 8-week, €100 build.", k.MUTE, False),
        ("Physical fluid rig", "Folded in", "Subtle and fragile — reborn as digital cymatics on screen.", k.AMBER, False),
        ("Poetry", "Dropped", "Negligible measurable body effect; it diluted the story.", k.MUTE, False),
        ("Bowl + heartbeat + body", "Chosen", "Research-backed, interactive, and every discipline contributes.", k.GREEN, True),
    ]
    y0, rh = 2.78, 0.66
    for i, (idea, status, why, col, keep) in enumerate(rows):
        ry = y0 + i * (rh + 0.08)
        c = k.card(s, M, ry, CW, rh, fill=(k.PANEL2 if keep else k.PANEL),
                   stroke=(k.GREEN if keep else k.HAIR), radius=0.10,
                   stroke_w=(1.4 if keep else 1.0))
        k.text(s, M + 0.34, ry, 3.4, rh, idea, font=k.DISPLAY, size=19,
               color=(k.INK if keep else k.BODY), anchor=MSO_ANCHOR.MIDDLE)
        pill = k.box(s, M + 3.9, ry + rh / 2 - 0.18, 1.35, 0.36, MSO_SHAPE.ROUNDED_RECTANGLE)
        pill.adjustments[0] = 0.5
        k.solid(pill, col, alpha=18 if not keep else 22)
        k.line(pill, col, 1.0)
        k.text(s, M + 3.9, ry + rh / 2 - 0.17, 1.35, 0.34, status, font=k.SANS, size=11,
               color=col, bold=True, tracking=1.4, align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE)
        k.text(s, M + 5.5, ry, CW - 5.8, rh, why, font=k.SANS, size=13,
               color=(k.BODY if keep else k.MUTE), anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        eff.append(E(c, "wipe", "after" if i == 0 else "after", 420, 30))
    k.footer(s, 6, section="What we discarded & why")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 07 · REQUIREMENTS
def s07():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    eff = head(s, "REQUIREMENTS & SCOPE", "What the piece must do")
    cols = [
        ("MUST", k.GREEN, ["React in real time", "Bowl sound → colour", "Heartbeat → pulse", "Run on one laptop"]),
        ("SHOULD", k.CYAN, ["Track the body", "Work fully offline", "Feel calm & immersive", "Each discipline visible"]),
        ("COULD", k.VIOLET, ["Extra instruments", "Background removal", "Recorded keepsake", "Multiple people"]),
        ("WON'T", k.MUTE, ["Medical claims", "Cloud / internet", "Paid services", "Custom hardware lab"]),
    ]
    cw, gx, x0, y0, ch = 2.85, 0.24, M, 2.30, 3.1
    for i, (label, col, items) in enumerate(cols):
        cx = x0 + i * (cw + gx)
        c = k.card(s, cx, y0, cw, ch, fill=k.PANEL, stroke=k.HAIR, radius=0.06)
        head_b = k.box(s, cx, y0, cw, 0.52, MSO_SHAPE.RECTANGLE)
        k.solid(head_b, col, alpha=16)
        k.line(head_b, col, 1.0)
        k.text(s, cx, y0 + 0.06, cw, 0.4, label, font=k.SANS, size=14, color=col,
               bold=True, tracking=3.0, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j, it in enumerate(items):
            k.rich(s, cx + 0.28, y0 + 0.76 + j * 0.54, cw - 0.5, 0.5,
                   [("—  ", col, True), (it, k.BODY, False)], size=12.5,
                   anchor=MSO_ANCHOR.MIDDLE)
        eff.append(E(c, "float", "after" if i == 0 else "with", 520, 60))
    strip = k.card(s, M, 5.7, CW, 0.62, fill="0F0F18", stroke=k.HAIR, radius=0.16)
    k.rich(s, M + 0.34, 5.7, CW - 0.7, 0.62,
           [("Hard constraints:   ", k.GOLD, True),
            ("€100 budget   ·   no internet at the demo   ·   one webcam   ·   six disciplines   ·   ready for 19 June", k.BODY, False)],
           size=13, anchor=MSO_ANCHOR.MIDDLE)
    eff.append(E(strip, "fade", "after", 480, 60))
    k.footer(s, 7, section="Requirements (MoSCoW)")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 08 · CONCEPT
def s08():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 6.7, 3.2, 5.2, k.VIOLET, 20)
    eff = head(s, "THE CONCEPT", "Three signals, one image")
    big = k.text(s, M, 2.2, CW, 1.3,
                 "One person. One bowl. One camera. — woven into a single living visual field.",
                 font=k.DISPLAY, italic=True, size=26, color=k.INK, align=PP_ALIGN.CENTER,
                 line_spacing=1.1)
    items = [
        ("THE BOWL", "sets the colour", "Its frequency picks the chakra hue and the cymatic pattern.", k.GREEN),
        ("THE HEARTBEAT", "sets the pulse", "Every beat makes the whole image gently throb — and tints it.", k.ROSE),
        ("THE BODY", "sets the shape", "The camera finds you, and the visuals wrap around your form.", k.CYAN),
    ]
    cw, gx, x0, y0, ch = 3.78, 0.30, M, 3.7, 2.3
    for i, (label, verb, body, col) in enumerate(items):
        cx = x0 + i * (cw + gx)
        c = k.card(s, cx, y0, cw, ch, fill=k.PANEL, stroke=k.HAIR, radius=0.07)
        dot = k.box(s, cx + cw / 2 - 0.12, y0 + 0.34, 0.24, 0.24, MSO_SHAPE.OVAL)
        k.solid(dot, col)
        k.text(s, cx, y0 + 0.72, cw, 0.34, label, font=k.SANS, size=12.5, color=col,
               bold=True, tracking=2.4, align=PP_ALIGN.CENTER)
        k.text(s, cx, y0 + 1.06, cw, 0.4, verb, font=k.DISPLAY, italic=True, size=20,
               color=k.INK, align=PP_ALIGN.CENTER)
        k.text(s, cx + 0.32, y0 + 1.56, cw - 0.64, 0.7, body, font=k.SANS, size=12.5,
               color=k.MUTE, align=PP_ALIGN.CENTER, line_spacing=1.18)
        eff.append(E(c, "float", "after" if i == 0 else "with", 540, 70))
        eff.append(E(dot, "zoom", "with", 360, 0))
    eff.insert(3, E(big, "fade", "after", 640, 80))
    k.footer(s, 8, section="The concept")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "morph", 900)


# ============================================================ 08b · USER-CENTRED
def s08b():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 11.3, 5.0, 4.2, k.CYAN, 16)
    eff = head(s, "USER-CENTRED DESIGN", "Anyone can walk up and feel it", accent=k.CYAN)
    pts = [
        ("No instructions", "You are the interface — step in and it simply responds. Nothing to read, nothing to operate.", k.CYAN),
        ("Made for any body", "MediaPipe tracks any height and any pose, so the visuals find every visitor the same way.", k.VIOLET),
        ("Calming by design", "Slow motion, soft colour and a dark, quiet space — built to settle a stressed visitor, not overload them.", k.GREEN),
        ("The space leads you", "One bowl, one focus, black surroundings — the piece guides the whole experience without a word.", k.GOLD),
    ]
    cw, ch, gx, gy = 5.5, 1.78, 0.4, 0.3
    x0, y0 = M, 2.28
    for i, (t, b, col) in enumerate(pts):
        cx = x0 + (i % 2) * (cw + gx)
        cy = y0 + (i // 2) * (ch + gy)
        c = k.card(s, cx, cy, cw, ch, fill=k.PANEL, stroke=k.HAIR, radius=0.06)
        tab = k.box(s, cx, cy + 0.18, 0.06, ch - 0.36, MSO_SHAPE.ROUNDED_RECTANGLE)
        k.solid(tab, col)
        k.text(s, cx + 0.34, cy + 0.26, cw - 0.6, 0.4, t, font=k.DISPLAY, size=21, color=k.INK)
        k.text(s, cx + 0.34, cy + 0.76, cw - 0.7, 0.9, b, font=k.SANS, size=13, color=k.MUTE,
               line_spacing=1.2)
        eff.append(E(c, "float", "after" if i == 0 else "with", 520, 60))
    note = k.text(s, M, 6.18, CW, 0.4,
                  "Every design choice met one test: would a stressed stranger understand it in five seconds?",
                  font=k.DISPLAY, italic=True, size=15, color=k.BODY, align=PP_ALIGN.CENTER)
    eff.append(E(note, "fade", "after", 460, 60))
    k.footer(s, section="User-centred design")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 09 · PIPELINE
def s09():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    eff = head(s, "HOW IT WORKS", "Inputs  →  processing  →  outputs")
    colx = [M + 0.1, 5.25, 9.4]
    colhead = [("INPUTS", k.CYAN), ("PROCESSING", k.VIOLET), ("OUTPUTS", k.GOLD)]
    for i, (lab, col) in enumerate(colhead):
        k.text(s, colx[i], 2.18, 3.3, 0.34, lab, font=k.SANS, size=12, color=col,
               bold=True, tracking=2.8)
    inputs = [("Bowl audio", "microphone", k.GREEN),
              ("Heartbeat", "Arduino pulse sensor", k.ROSE),
              ("The body", "webcam", k.CYAN)]
    procs = [("FFT  /  Goertzel", "nearest chakra, 396–963 Hz", k.VIOLET),
             ("BPM  →  LFO", "smoothed pulse", k.VIOLET),
             ("MediaPipe pose", "head, hands, torso", k.VIOLET)]
    outs = [("Colour & cymatics", "", k.GREEN),
            ("Pulsing throb & tint", "", k.ROSE),
            ("Aura wraps the body", "", k.CYAN)]

    def chip(x, y, title, sub, col, w=3.4):
        c = k.card(s, x, y, w, 0.92, fill=k.PANEL, stroke=k.HAIR, radius=0.12)
        tab = k.box(s, x, y + 0.16, 0.06, 0.60, MSO_SHAPE.ROUNDED_RECTANGLE)
        k.solid(tab, col)
        k.text(s, x + 0.26, y + (0.16 if sub else 0.30), w - 0.4, 0.4, title,
               font=k.SANS, size=14.5, color=k.INK, bold=True)
        if sub:
            k.text(s, x + 0.26, y + 0.52, w - 0.4, 0.34, sub, font=k.SANS, size=11.5, color=k.MUTE)
        return c

    ys = [2.66, 3.92, 5.18]
    for r in range(3):
        ci = chip(colx[0], ys[r], inputs[r][0], inputs[r][1], inputs[r][2])
        cp = chip(colx[1], ys[r], procs[r][0], procs[r][1], procs[r][2])
        co = chip(colx[2], ys[r], outs[r][0], outs[r][1], outs[r][2])
        a1 = k.text(s, colx[0] + 3.45, ys[r], 0.55, 0.92, "→", font=k.SANS, size=26,
                    color=k.HAIR, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        a2 = k.text(s, colx[1] + 3.45, ys[r], 0.55, 0.92, "→", font=k.SANS, size=26,
                    color=k.HAIR, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        eff += [E(ci, "fade", "after" if r == 0 else "after", 360, 20),
                E(a1, "fade", "with", 260, 0), E(cp, "fade", "with", 360, 0),
                E(a2, "fade", "with", 260, 0), E(co, "fade", "with", 360, 0)]
    note = k.text(s, M, 6.5, CW, 0.4,
                  "All of it runs on the GPU in TouchDesigner — one laptop, one webcam, no internet.",
                  font=k.DISPLAY, italic=True, size=14.5, color=k.BODY, align=PP_ALIGN.CENTER)
    eff.append(E(note, "fade", "after", 460, 60))
    k.footer(s, 9, section="System pipeline")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 10 · BOWL MEASURED
def s10():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    eff = head(s, "PROTOTYPE & FEASIBILITY", "A working prototype — and we tested it")
    specs = [("bowl_p4_x36_794x282.png", "EMPTY BOWL", "647 Hz", k.GREEN),
             ("bowl_p6_x47_805x282.png", "+ 0.1 L WATER", "629 Hz", k.CYAN)]
    iw = 5.2
    x0, y0 = M, 2.18
    for i, (fn, lab, val, col) in enumerate(specs):
        cx = x0 + i * (iw + 0.4)
        c = k.card(s, cx - 0.12, y0 - 0.12, iw + 0.24, 2.06, fill="0A0A12", stroke=k.HAIR, radius=0.05)
        p = s.shapes.add_picture(img(fn), Inches(cx), Inches(y0), width=Inches(iw))
        p.crop_top = 0.13; p.crop_bottom = 0.05; p.crop_left = 0.02; p.crop_right = 0.07
        k.text(s, cx + 0.04, y0 + 1.6, iw, 0.3, lab, font=k.SANS, size=10.5, color=col,
               bold=True, tracking=1.6)
        k.text(s, cx + iw - 2.0, y0 + 1.56, 2.0, 0.34, val, font=k.DISPLAY, italic=True,
               size=16, color=k.INK, align=PP_ALIGN.RIGHT)
        eff += [E(c, "fade", "after" if i == 0 else "with", 420, 40),
                E(p, "fade", "with", 600, 0)]
    cap = k.text(s, M, y0 + 1.98, CW, 0.3,
                 "Measured live on the real bowl — f ∝ √(k/m): more water → more mass → a lower note. Overtones near 1756 Hz make it an inharmonic resonator.",
                 font=k.SANS, size=11, color=k.MUTE, italic=True)
    eff.append(E(cap, "fade", "after", 360, 40))
    vy = 4.66
    vc = k.card(s, M, vy, CW, 1.86, fill=k.PANEL, stroke=k.GREEN, radius=0.05, stroke_w=1.0)
    k.text(s, M + 0.34, vy + 0.2, 8, 0.3, "ALREADY VALIDATED", font=k.SANS, size=11,
           color=k.GREEN, bold=True, tracking=2.2)
    checks = ["Bowl frequencies measured — 647 Hz, falling to 629 Hz with water",
              "Body tracking runs live — head & torso at 0.97 / 0.99 confidence, no browser",
              "Particle system smoke-tested — 1031 of 2048 gather to a still hand",
              "The whole show runs on one laptop, fully offline"]
    cw2 = (CW - 0.7) / 2
    for i, ct in enumerate(checks):
        ckx = M + 0.34 + (i % 2) * cw2
        cky = vy + 0.64 + (i // 2) * 0.56
        k.rich(s, ckx, cky, cw2 - 0.2, 0.5, [("✓  ", k.GREEN, True), (ct, k.BODY, False)],
               size=12.5, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    eff.append(E(vc, "float", "after", 520, 70))
    k.footer(s, section="Prototype & feasibility")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 11 · SOUND -> COLOUR
def s11():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    eff = head(s, "SOUND BECOMES COLOUR", "Frequency → chakra → cymatics")
    # chakra list (left)
    lx, ly = M, 2.3
    k.text(s, lx, ly - 0.04, 5.6, 0.3, "NEAREST SOLFEGGIO CHAKRA  (396–963 Hz)", font=k.SANS,
           size=10.5, color=k.MUTE, bold=True, tracking=1.8)
    rows_eff = []
    for i, (name, hz, hx, kw) in enumerate(k.CHAKRAS):
        ry = ly + 0.42 + i * 0.46
        is_heart = (name == "Heart")
        if is_heart:
            hb = k.card(s, lx - 0.1, ry - 0.04, 5.75, 0.44, fill=hx, stroke=None, radius=0.2, fill_alpha=14)
            rows_eff.append(E(hb, "fade", "with", 300, 0))
        sw = k.box(s, lx, ry, 0.30, 0.30, MSO_SHAPE.OVAL)
        k.solid(sw, hx)
        k.text(s, lx + 0.46, ry - 0.04, 1.9, 0.34, name, font=k.SANS, size=13,
               color=(k.INK if is_heart else k.BODY), bold=is_heart, anchor=MSO_ANCHOR.MIDDLE)
        k.text(s, lx + 2.3, ry - 0.04, 1.0, 0.34, "%d Hz" % hz, font=k.SANS, size=12.5,
               color=k.MUTE, anchor=MSO_ANCHOR.MIDDLE)
        if is_heart:
            k.text(s, lx + 3.4, ry - 0.04, 2.3, 0.34, "← the bowl, 647 Hz", font=k.DISPLAY,
                   italic=True, size=13, color=hx, anchor=MSO_ANCHOR.MIDDLE)
        else:
            k.text(s, lx + 3.4, ry - 0.04, 2.2, 0.34, kw, font=k.SANS, size=11.5,
                   color=k.MUTE, anchor=MSO_ANCHOR.MIDDLE, italic=True)
        rows_eff.append(E(sw, "wipe", "after" if i == 0 else "with", 300, 0))
    eff += rows_eff
    # cymatics montage (right) 2x2
    cyms = [("cymatic_p2_x2_997x952.png", "Root"), ("cymatic_p2_x3_996x951.png", "Heart"),
            ("cymatic_p2_x4_993x947.png", "Throat"), ("cymatic_p2_x5_990x940.png", "Third Eye")]
    gx0, gy0, isz, gp = 7.25, 2.4, 2.46, 0.18
    for i, (fn, lab) in enumerate(cyms):
        px = gx0 + (i % 2) * (isz + gp)
        py = gy0 + (i // 2) * (isz + gp)
        p = s.shapes.add_picture(img(fn), Inches(px), Inches(py), width=Inches(isz), height=Inches(isz))
        p.crop_top = 0.12; p.crop_left = 0.03; p.crop_right = 0.03; p.crop_bottom = 0.02
        eff.append(E(p, "zoom", "after" if i == 0 else "with", 520, 40))
    cap = k.text(s, gx0, gy0 + 2 * isz + gp + 0.06, 2 * isz + gp, 0.4,
                 "One equation, sin(kx)·sin(ky) — denser as the note rises.", font=k.SANS,
                 size=11.5, color=k.MUTE, italic=True, align=PP_ALIGN.CENTER)
    eff.append(E(cap, "fade", "after", 420, 40))
    k.footer(s, 11, section="Sound → colour")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 12 · HEARTBEAT
def s12():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 3.0, 3.4, 4.0, k.ROSE, 18)
    eff = head(s, "THE HEARTBEAT", "Your pulse sets the rhythm — and the colour", accent=k.ROSE)
    body = k.text(s, M, 2.36, 5.2, 2.4,
                  "An Arduino pulse sensor reads your heart rate, then smooths it into a low-frequency pulse that drives the whole image.\n\nEvery beat, the visuals gently throb — the way a calm body breathes.",
                  font=k.SANS, size=15.5, color=k.BODY, line_spacing=1.26)
    chain = k.rich(s, M, 4.7, 5.4, 0.5,
                   [("pulse sensor  ", k.ROSE, True), ("→  BPM  →  smoothed  →  ", k.MUTE, False),
                    ("heartbeat LFO", k.ROSE, True)], size=13.5)
    eff += [E(body, "fade", "after", 620, 80), E(chain, "fade", "after", 460, 60)]
    # BPM colour scale (right)
    sx, sy, sw = 6.7, 2.7, 5.6
    k.text(s, sx, sy - 0.4, sw, 0.34, "BPM  →  COLOUR", font=k.SANS, size=11,
           color=k.MUTE, bold=True, tracking=2.4)
    bands = [("VIOLET", "< 62", "calm", k.VIOLET), ("CYAN", "62–72", "resting", k.CYAN),
             ("AMBER", "72–85", "neutral", k.AMBER), ("ORANGE", "≥ 85", "elevated", k.ORANGE)]
    bw = (sw - 3 * 0.16) / 4
    for i, (nm, rng, mood, col) in enumerate(bands):
        bx = sx + i * (bw + 0.16)
        bar = k.box(s, bx, sy, bw, 0.7, MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.adjustments[0] = 0.18
        k.solid(bar, col)
        k.text(s, bx, sy + 0.84, bw, 0.34, rng, font=k.DISPLAY, size=18, color=k.INK,
               align=PP_ALIGN.CENTER)
        k.text(s, bx, sy + 1.24, bw, 0.3, mood, font=k.SANS, size=11, color=k.MUTE,
               align=PP_ALIGN.CENTER, tracking=1.0)
        eff.append(E(bar, "wipe", "after" if i == 0 else "with", 360, 0))
    legend = k.text(s, sx, sy + 1.8, sw, 0.5,
                    "Calmer heart → cool violets and cyans.   Faster heart → warm ambers and oranges.",
                    font=k.SANS, size=12.5, color=k.BODY, line_spacing=1.2)
    eff.append(E(legend, "fade", "after", 460, 60))
    k.footer(s, 12, section="Heartbeat → pulse")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 13 · BODY
def s13():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 10.4, 3.6, 4.4, k.CYAN, 18)
    eff = head(s, "THE BODY", "The visuals wrap around you", accent=k.CYAN)
    body = k.text(s, M, 2.36, 6.0, 2.2,
                  "A single webcam feeds MediaPipe pose tracking — it finds your head, hands and torso many times a second, fully on-device.\n\nThe aura and a 2,048-particle field follow you: particles gather to still hands, and scatter from fast ones, so motion becomes light.",
                  font=k.SANS, size=15.5, color=k.BODY, line_spacing=1.26)
    feats = [("On-device tracking", "no internet, no markers", k.CYAN),
             ("2,048 GPU particles", "gather & scatter with your hands", k.VIOLET),
             ("One-Euro smoothing", "the aura glides, never snaps", k.GREEN)]
    fy = 4.78
    for i, (t, sub, col) in enumerate(feats):
        c = k.card(s, M, fy + i * 0.62, 6.0, 0.54, fill=k.PANEL, stroke=k.HAIR, radius=0.14)
        d = k.box(s, M + 0.22, fy + i * 0.62 + 0.17, 0.2, 0.2, MSO_SHAPE.OVAL)
        k.solid(d, col)
        k.rich(s, M + 0.6, fy + i * 0.62, 5.3, 0.54,
               [(t + "  —  ", k.INK, True), (sub, k.MUTE, False)], size=12.5,
               anchor=MSO_ANCHOR.MIDDLE)
        eff.append(E(c, "wipe", "after", 380, 30))
    # stylised body-with-aura on the right
    ax, ay = 9.9, 3.3
    rings = [(2.0, k.CYAN, 10), (1.5, k.VIOLET, 16), (1.0, k.GOLD, 22)]
    ring_shapes = []
    for r, col, al in rings:
        g = k.glow(s, ax, ay, r, col, al)
        ring_shapes.append(g)
    figure = k.box(s, ax - 0.22, ay - 0.9, 0.44, 1.8, MSO_SHAPE.ROUNDED_RECTANGLE)
    figure.adjustments[0] = 0.5
    k.solid(figure, k.INK, alpha=85)
    headd = k.box(s, ax - 0.22, ay - 1.3, 0.44, 0.44, MSO_SHAPE.OVAL)
    k.solid(headd, k.INK, alpha=85)
    for g in ring_shapes:
        eff.append(E(g, "fade", "with", 700, 0))
    eff += [E(figure, "fade", "after", 460, 40), E(headd, "fade", "with", 460, 0)]
    k.footer(s, 13, section="Body → shape")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 700)


# ============================================================ 14 · EXPERIENCE / DEMO
def s14():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    eff = head(s, "THE EXPERIENCE", "You become part of the artwork")
    steps = ["You step in front of the installation.",
             "The Tibetan bowl is played.",
             "Your heartbeat is read as the bowl keeps ringing.",
             "Sound, pulse and body turn into a live visual field.",
             "As you settle, the colours cool and slow — you see yourself calm."]
    y0 = 2.34
    for i, st in enumerate(steps):
        ry = y0 + i * 0.62
        num = k.box(s, M, ry, 0.42, 0.42, MSO_SHAPE.OVAL)
        k.solid(num, k.PANEL2); k.line(num, k.GOLD, 1.2)
        k.text(s, M, ry + 0.02, 0.42, 0.38, str(i + 1), font=k.DISPLAY, size=16,
               color=k.GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = k.text(s, M + 0.62, ry, 5.5, 0.5,
                    st, font=k.SANS, size=13.5, color=(k.INK if i == 4 else k.BODY),
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        eff += [E(num, "zoom", "after" if i == 0 else "after", 360, 20),
                E(tx, "fade", "with", 360, 0)]
    # TD screenshot on the right
    px, py, pw = 7.0, 2.5, 5.3
    c = k.card(s, px - 0.12, py - 0.12, pw + 0.24, pw * 668 / 1259 + 0.24, fill="0A0A12",
               stroke=k.HAIR, radius=0.05)
    p = s.shapes.add_picture(img("tdnet_p4_x34_1259x668.png"), Inches(px), Inches(py), width=Inches(pw))
    cap = k.text(s, px, py + pw * 668 / 1259 + 0.18, pw, 0.34,
                 "The real engine: our TouchDesigner network, running live.", font=k.SANS,
                 size=11.5, color=k.MUTE, italic=True, align=PP_ALIGN.CENTER)
    demo = k.box(s, px + pw - 2.3, py + 0.18, 2.1, 0.5, MSO_SHAPE.ROUNDED_RECTANGLE)
    demo.adjustments[0] = 0.5
    k.solid(demo, k.GOLD)
    k.text(s, px + pw - 2.3, py + 0.18, 2.1, 0.5, "▶  LIVE DEMO", font=k.SANS, size=12.5,
           color="0A0A12", bold=True, tracking=1.6, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    eff += [E(c, "fade", "after", 460, 60), E(p, "fade", "with", 640, 0),
            E(demo, "zoom", "after", 460, 60), E(cap, "fade", "with", 360, 0)]
    k.footer(s, 14, section="The experience")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "morph", 900)


# ============================================================ 15 · REFLECTION
def s15():
    s = prs.slides.add_slide(BLANK)
    k.background(s)
    k.glow(s, 6.7, 3.0, 5.0, k.VIOLET, 16)
    eff = head(s, "REFLECTION & NEXT STEPS", "What we fixed — and where we go next")
    # left: fixed
    lc = k.card(s, M, 2.34, 5.5, 3.1, fill=k.PANEL, stroke=k.HAIR, radius=0.06)
    k.text(s, M + 0.34, 2.56, 5.0, 0.34, "FROM THE MIDTERM FEEDBACK", font=k.SANS, size=11,
           color=k.ROSE, bold=True, tracking=2.0)
    fixes = [("“Idea wasn't visualised”", "now a live, reacting installation"),
             ("“Too complex”", "one clear story: sound + body + calm"),
             ("“Not explained clearly”", "this deck and a live demo")]
    for i, (a, b) in enumerate(fixes):
        yy = 3.02 + i * 0.74
        k.rich(s, M + 0.34, yy, 5.0, 0.4, [(a, k.BODY, False)], size=13.5)
        k.rich(s, M + 0.34, yy + 0.32, 5.0, 0.4, [("→  ", k.GREEN, True), (b, k.INK, True)], size=13)
    # right: next
    rc = k.card(s, 6.92, 2.34, 5.5, 3.1, fill=k.PANEL, stroke=k.HAIR, radius=0.06)
    k.text(s, 7.26, 2.56, 5.0, 0.34, "WHAT'S NEXT", font=k.SANS, size=11, color=k.CYAN,
           bold=True, tracking=2.0)
    nexts = [("More instruments", "kalimba, tuning fork, chimes — to reach every chakra"),
             ("Live Arduino BPM", "flash the sensor, verify the pulse on stage"),
             ("Floating body", "drop the room so the person floats on black")]
    for i, (a, b) in enumerate(nexts):
        yy = 3.02 + i * 0.74
        k.rich(s, 7.26, yy, 5.0, 0.4, [("—  ", k.CYAN, True), (a, k.INK, True)], size=13.5)
        k.text(s, 7.5, yy + 0.32, 4.8, 0.4, b, font=k.SANS, size=12, color=k.MUTE, line_spacing=1.1)
    eff += [E(lc, "float", "after", 540, 70), E(rc, "float", "with", 540, 0)]
    bar = k.chakra_bar(s, M, 5.84, 5.0, h=0.14)
    ty = k.text(s, 6.5, 5.7, CW - 5.6, 0.5, "Thank you.", font=k.DISPLAY, italic=True,
                size=28, color=k.INK, align=PP_ALIGN.RIGHT)
    nm = k.text(s, 6.5, 6.28, CW - 5.6, 0.3, "Creative Bowl Lab  ·  Group 5  ·  TU/e", font=k.SANS,
                size=10.5, color=k.MUTE, tracking=1.6, align=PP_ALIGN.RIGHT)
    eff += [E(seg, "wipe", "after" if i == 0 else "with", 320, 0) for i, seg in enumerate(bar)]
    eff += [E(ty, "fade", "after", 560, 60), E(nm, "fade", "after", 360, 30)]
    k.footer(s, 15, section="Reflection")
    anim.add_entrance(s, eff)
    anim.add_transition(s, "fade", 800)


for fn in (s01, s02, s03, s04, s05, s06, s07, s08, s08b, s09, s10, s11, s12, s13, s14, s15):
    fn()

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
