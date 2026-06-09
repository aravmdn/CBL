"""Design-system primitives for the Creative Bowl Lab deck.

A small toolkit over python-pptx: gradient backgrounds, soft radial glows,
typographic helpers (Garamond display + Segoe UI sans), rounded "cards", pills,
hairlines, the recurring chakra-spectrum motif, and footers. Everything returns
the created shape so the caller can attach entrance animations.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ---- palette ---------------------------------------------------------------
BG0   = "06060C"   # deepest background
BG1   = "0C0C16"   # background gradient bottom
BG2   = "10101C"
PANEL = "14141F"   # card fill
PANEL2= "1A1A2A"   # raised card
HAIR  = "2C2C40"   # hairlines / strokes
INK   = "F3F2F8"   # titles
BODY  = "C9C8D8"   # body text
MUTE  = "8A899F"   # captions / muted

VIOLET = "8B5CF6"
CYAN   = "22D3EE"
AMBER  = "F5A623"
ORANGE = "FB7138"
ROSE   = "FB7185"
GREEN  = "34D399"
GOLD   = "D9B775"

# chakra: (name, Hz, hex, keyword)
CHAKRAS = [
    ("Root",        396, "F2555A", "release"),
    ("Sacral",      417, "FB923C", "change"),
    ("Solar Plexus",528, "FACC15", "energy"),
    ("Heart",       639, "34D399", "connection"),
    ("Throat",      741, "38BDF8", "expression"),
    ("Third Eye",   852, "818CF8", "intuition"),
    ("Crown",       963, "C084FC", "awakening"),
]

DISPLAY = "Garamond"
SANS    = "Segoe UI"

EMU_IN = 914400
SW, SH = 13.333, 7.5
MARGIN = 0.92
CW = SW - 2 * MARGIN   # content width


def rgb(h):
    return RGBColor.from_string(h)


# ---- low-level fill / line -------------------------------------------------
def _spPr(shape):
    return shape._element.spPr


def _clear_fill(spPr):
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
              "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(t))
        if e is not None:
            spPr.remove(e)


def _insert_fill(spPr, el):
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(el)
    else:
        spPr.append(el)


def no_shadow(shape):
    shape.shadow.inherit = False


def no_line(shape):
    shape.line.fill.background()


def line(shape, hexcolor, w=1.0, alpha=None):
    shape.line.color.rgb = rgb(hexcolor)
    shape.line.width = Pt(w)
    if alpha is not None:
        srgb = shape.line._get_or_add_ln().find(qn("a:solidFill")).find(qn("a:srgbClr"))
        srgb.append(parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/'
                              'drawingml/2006/main" val="%d"/>' % int(alpha * 1000)))


def solid(shape, hexcolor, alpha=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hexcolor)
    if alpha is not None:
        srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
        srgb.append(parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/'
                              'drawingml/2006/main" val="%d"/>' % int(alpha * 1000)))


def gradient(shape, top_hex, bottom_hex, angle_deg=90):
    spPr = _spPr(shape)
    _clear_fill(spPr)
    ang = int(angle_deg * 60000)
    xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="%s"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="%s"/></a:gs>'
        '</a:gsLst><a:lin ang="%d" scaled="1"/></a:gradFill>' % (top_hex, bottom_hex, ang)
    )
    _insert_fill(spPr, parse_xml(xml))


def radial(shape, hexcolor, inner_alpha=55, outer_alpha=0):
    """Soft circular glow: inner colour fades to transparent at the edge."""
    spPr = _spPr(shape)
    _clear_fill(spPr)
    no_line(shape)
    xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '</a:gsLst><a:path path="circle">'
        '<a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path></a:gradFill>'
        % (hexcolor, int(inner_alpha * 1000), hexcolor, int(outer_alpha * 1000))
    )
    _insert_fill(spPr, parse_xml(xml))


# ---- shapes ----------------------------------------------------------------
def box(slide, l, t, w, h, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    no_shadow(sp)
    no_line(sp)
    return sp


def card(slide, l, t, w, h, fill=PANEL, stroke=HAIR, radius=0.06, stroke_w=1.0,
         fill_alpha=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    no_shadow(sp)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    solid(sp, fill, fill_alpha)
    if stroke:
        line(sp, stroke, stroke_w)
    else:
        no_line(sp)
    return sp


# ---- background ------------------------------------------------------------
def background(slide, top=BG1, bottom=BG0):
    # Flat near-black: near-identical dark gradients band badly when rendered,
    # so we use a solid base and rely on the radial glows for depth.
    bg = box(slide, -0.06, -0.06, SW + 0.12, SH + 0.12)
    solid(bg, "07070E")
    return bg


def glow(slide, cx, cy, r, hexcolor, inner_alpha=42):
    g = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(cx - r), Inches(cy - r),
                               Inches(2 * r), Inches(2 * r))
    no_shadow(g)
    radial(g, hexcolor, inner_alpha=inner_alpha)
    return g


# ---- text ------------------------------------------------------------------
def _style_run(r, font, size, color, bold, italic, tracking):
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)
    if tracking is not None:
        r._r.get_or_add_rPr().set("spc", str(int(tracking * 100)))


def text(slide, l, t, w, h, s, font=SANS, size=18, color=BODY, bold=False,
         italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         tracking=None, line_spacing=1.06, space_after=0, wrap=True,
         caps=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = s.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = ln.upper() if caps else ln
        _style_run(r, font, size, color, bold, italic, tracking)
    return tb


def rich(slide, l, t, w, h, segments, size=18, font=SANS, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.1):
    """segments: list of (text, color, bold) on one paragraph."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for seg in segments:
        txt, color, bold = seg
        r = p.add_run()
        r.text = txt
        _style_run(r, font, size, color, bold, False, None)
    return tb


def kicker(slide, l, t, s, color=GOLD, size=12.5, w=None):
    return text(slide, l, t, w or 8, 0.34, s, font=SANS, size=size, color=color,
                bold=True, tracking=3.2, caps=True)


def rule(slide, l, t, w, color=GOLD, weight=1.6, alpha=None):
    ln = slide.shapes.add_connector(2, Inches(l), Inches(t), Inches(l + w), Inches(t))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(weight)
    no_shadow(ln)
    if alpha is not None:
        srgb = ln.line._get_or_add_ln().find(qn("a:solidFill")).find(qn("a:srgbClr"))
        srgb.append(parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/'
                              'drawingml/2006/main" val="%d"/>' % int(alpha * 1000)))
    return ln


# ---- recurring motifs ------------------------------------------------------
def chakra_bar(slide, l, t, w, h=0.12, gap=0.04, alpha=None):
    """Seven-segment chakra spectrum — the deck's signature motif. Returns the
    list of segment shapes (so they can wipe in one by one)."""
    n = len(CHAKRAS)
    seg_w = (w - gap * (n - 1)) / n
    segs = []
    for i, (_, _, hx, _) in enumerate(CHAKRAS):
        x = l + i * (seg_w + gap)
        s = box(slide, x, t, seg_w, h, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        try:
            s.adjustments[0] = 0.5
        except Exception:
            pass
        solid(s, hx, alpha)
        segs.append(s)
    return segs


_FOOTN = [1]  # auto slide counter (title is slide 1, has no footer)


def footer(slide, idx=None, total=16, section=""):
    _FOOTN[0] += 1
    n = _FOOTN[0]
    text(slide, MARGIN, SH - 0.5, 5, 0.3, "CREATIVE BOWL LAB", font=SANS,
         size=9, color=MUTE, bold=True, tracking=2.4)
    text(slide, SW / 2 - 3, SH - 0.5, 6, 0.3, section, font=SANS, size=9,
         color=MUTE, tracking=2.4, align=PP_ALIGN.CENTER, caps=True)
    text(slide, SW - MARGIN - 3, SH - 0.5, 3, 0.3,
         "GROUP 5  ·  %02d / %02d" % (n, total), font=SANS, size=9,
         color=MUTE, bold=True, tracking=2.0, align=PP_ALIGN.RIGHT)


def header(slide, kicker_txt, title_txt, accent=GOLD, title_size=34):
    """Standard slide header: kicker + big Garamond title + accent rule.
    Returns (kicker_shape, title_shape, rule_shape)."""
    k = kicker(slide, MARGIN, 0.62, kicker_txt, color=accent)
    ttl = text(slide, MARGIN, 0.92, CW, 0.9, title_txt, font=DISPLAY,
               size=title_size, color=INK, bold=False, line_spacing=1.0)
    rl = rule(slide, MARGIN, 1.74, 0.86, color=accent, weight=2.2)
    return k, ttl, rl
